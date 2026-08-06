import os, subprocess, tempfile, zipfile, shutil, logging, io
import re as _re
import csv
import uuid

log = logging.getLogger(__name__)

import fitz
from PIL import Image

try:
    import win32com.client # type: ignore
    import pythoncom # type: ignore
    _TEM_WIN32 = True
except ImportError:
    _TEM_WIN32 = False

try:
    # pyrefly: ignore [missing-import]
    from pillow_heif import register_heif_opener
    register_heif_opener()
except ImportError:
    pass


# ─── Motor ───────────────────────────────────────────────────

def _encontrar_soffice():
    if shutil.which("soffice"):     return "soffice"
    if shutil.which("libreoffice"): return "libreoffice"
    for c in [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        r"C:\Program Files\LibreOffice 7\program\soffice.exe",
        r"C:\Program Files\LibreOffice 24\program\soffice.exe",
    ]:
        if os.path.exists(c): return c
    return None


SOFFICE = _encontrar_soffice()

if _TEM_WIN32:    MOTOR = "office"
elif SOFFICE:     MOTOR = "libreoffice"
else:             MOTOR = None

log.info(f"[Prisma] Motor: {MOTOR or 'NENHUM'}")


# ─── Conversões disponíveis ───────────────────────────────────

CONVERSOES = {
    "csv":  ["xlsx", "pdf", "png", "jpg", "docx", "pptx", "json", "webp", "heic"],
    "xlsx": ["csv",  "pdf", "png", "jpg", "docx", "pptx", "json", "webp", "heic"],
    "json": ["csv",  "xlsx", "pdf", "png", "jpg", "webp", "heic"],
    "pdf":  ["docx", "pptx", "ppt", "png", "jpg", "webp", "heic", "xlsx", "csv", "txt"],
    "docx": ["pdf",  "png", "jpg", "webp", "heic", "xlsx", "csv", "pptx"],
    "ppt":  ["pdf",  "docx", "xlsx", "csv", "png", "jpg", "webp", "heic", "pptx"],
    "pptx": ["pdf",  "docx", "xlsx", "csv", "png", "jpg", "webp", "heic", "ppt"],
    "png":  ["pdf",  "jpg", "webp", "heic", "docx", "pptx", "txt"],
    "jpg":  ["pdf",  "png", "webp", "heic", "docx", "pptx", "txt"],
    "webp": ["pdf",  "png", "jpg", "heic", "docx", "pptx", "txt"],
    "heic": ["pdf",  "png", "jpg", "webp", "docx", "pptx", "txt"]
}

def obter_conversoes(extensao):
    ext = extensao.lower()
    if ext == "jpeg": ext = "jpg"
    return CONVERSOES.get(ext, [])
def obter_motor():               return MOTOR


# ─── Encoding ────────────────────────────────────────────────

def detectar_encoding(caminho: str) -> str:
    # PERF-007: lê arquivo uma única vez (BOM + chardet juntos)
    try:
        with open(caminho, "rb") as f:
            raw = f.read(10004)  # 4 bytes BOM + 10000 bytes para chardet
    except OSError:
        return "utf-8"

    # Detecta BOM
    if raw[:3] == b"\xef\xbb\xbf":                            return "utf-8-sig"
    if raw[:2] in (b"\xff\xfe", b"\xfe\xff"):               return "utf-16"

    # Usa chardet sobre os primeiros 10000 bytes
    try:
        import chardet
        res = chardet.detect(raw[:10000])
        if res and res.get("encoding"):
            return res["encoding"]
    except ImportError:
        pass

    # Fallback: tenta decodificar com encodings comuns
    for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
        try:
            raw[:1000].decode(enc)
            return enc
        except (UnicodeDecodeError, LookupError):
            continue
    return "utf-8"


# ─── LibreOffice headless ─────────────────────────────────────

def _soffice_convert(entrada, saida, fmt, timeout=120):
    if not SOFFICE: raise RuntimeError("LibreOffice não encontrado.")
    saida_dir  = os.path.dirname(os.path.abspath(saida))
    nome_base  = os.path.splitext(os.path.basename(entrada))[0]
    saida_auto = os.path.join(saida_dir, f"{nome_base}.{fmt}")
    r = subprocess.run(
        [SOFFICE, "--headless", "--norestore",
         "--convert-to", fmt, "--outdir", saida_dir,
         os.path.abspath(entrada)],
        capture_output=True, timeout=timeout
    )
    if r.returncode != 0:
        raise RuntimeError(f"LibreOffice: {r.stderr.decode('utf-8', errors='replace')}")
    if os.path.abspath(saida_auto) != os.path.abspath(saida) and os.path.exists(saida_auto):
        os.replace(saida_auto, saida)
    if not os.path.exists(saida):
        raise RuntimeError("LibreOffice não gerou o arquivo de saída.")


# ─── Pré-formata XLSX para PDF organizado ────────────────────

def _preparar_xlsx_para_pdf(caminho: str, orientacao: str = "retrato"):
    from openpyxl import load_workbook
    from openpyxl.worksheet.page import PageMargins
    wb = load_workbook(caminho)
    for ws in wb.worksheets:
        ws.page_setup.orientation = "landscape" if orientacao == "paisagem" else "portrait"
        ws.page_setup.fitToPage   = True
        ws.page_setup.fitToWidth  = 1
        ws.page_setup.fitToHeight = 0
        ws.page_margins = PageMargins(
            left=0.4, right=0.4, top=0.5, bottom=0.5, header=0.2, footer=0.2
        )
        if ws.max_row and ws.max_row > 1:
            ws.print_title_rows = "1:1"
        for col in ws.columns:
            letra   = col[0].column_letter
            max_len = max((len(str(c.value)) if c.value is not None else 0 for c in col), default=0)
            ws.column_dimensions[letra].width = min(max(max_len + 2, 8), 35)
    wb.save(caminho)


# ─── CSV ↔ XLSX ───────────────────────────────────────────────


def _carregar_csv(caminho, enc):
    try: import pandas as pd
    except ImportError: raise RuntimeError("pandas não instalado")
    try:
        with open(caminho, 'r', encoding=enc, errors='replace') as f:
            amostra = f.read(4096)
            f.seek(0)
            separador = csv.Sniffer().sniff(amostra).delimiter
    except Exception:
        separador = ','
    return pd.read_csv(caminho, encoding=enc, sep=separador, engine="c", encoding_errors="replace", on_bad_lines="skip", low_memory=False)

def csv_para_xlsx(entrada, saida):
    enc = detectar_encoding(entrada)
    df  = _carregar_csv(entrada, enc)
    df.fillna("").to_excel(saida, index=False, engine="openpyxl")

def xlsx_para_csv(entrada, saida):
    try: import pandas as pd
    except ImportError: raise RuntimeError("pandas não instalado")
    df = pd.read_excel(entrada, engine="openpyxl")
    df.fillna("").to_csv(saida, index=False, encoding="utf-8-sig")


# ─── XLSX → PDF ───────────────────────────────────────────────

def xlsx_para_pdf(entrada, saida, orientacao="retrato"):
    if MOTOR == "office":
        _office_xlsx_para_pdf(entrada, saida, orientacao)
    elif MOTOR == "libreoffice":
        _preparar_xlsx_para_pdf(entrada, orientacao)
        _soffice_convert(entrada, saida, "pdf")
    else:
        raise RuntimeError("Nenhum motor disponível.")

def _office_xlsx_para_pdf(entrada, saida, orientacao="retrato"):
    pythoncom.CoInitialize()
    excel = wb = None
    try:
        excel = win32com.client.Dispatch("Excel.Application"); excel.Visible = False
        wb    = excel.Workbooks.Open(os.path.abspath(entrada))
        for sheet in wb.Worksheets:
            ps = sheet.PageSetup
            ps.Orientation    = 2 if orientacao == "paisagem" else 1
            ps.Zoom           = False
            ps.FitToPagesWide = 1
            ps.FitToPagesTall = False
            ps.LeftMargin     = excel.InchesToPoints(0.4)
            ps.RightMargin    = excel.InchesToPoints(0.4)
            ps.TopMargin      = excel.InchesToPoints(0.5)
            ps.BottomMargin   = excel.InchesToPoints(0.5)
            ps.PrintTitleRows = "$1:$1"
        wb.ExportAsFixedFormat(0, os.path.abspath(saida))
    finally:
        if wb:    
            try: wb.Close(False)
            except Exception: pass
        if excel: 
            try: excel.Quit()
            except Exception: pass
        pythoncom.CoUninitialize()


# ─── CSV → PDF ────────────────────────────────────────────────

def csv_para_pdf(entrada, saida, orientacao="retrato"):
    enc  = detectar_encoding(entrada)
    temp = os.path.join(os.path.dirname(saida) or ".", f"_tmp_{uuid.uuid4().hex}.xlsx")
    try:
        df = _carregar_csv(entrada, enc)
        df.fillna("").to_excel(temp, index=False, engine="openpyxl")
        xlsx_para_pdf(temp, saida, orientacao=orientacao)
    finally:
        if os.path.exists(temp):
            try: os.remove(temp)
            except Exception: pass


# ─── DOCX → PDF ───────────────────────────────────────────────

def docx_para_pdf(entrada, saida, orientacao="retrato"):
    if MOTOR == "office":        _office_docx_para_pdf(entrada, saida)
    elif MOTOR == "libreoffice": _soffice_convert(entrada, saida, "pdf")
    else: raise RuntimeError("Nenhum motor disponível.")

def _office_docx_para_pdf(entrada, saida):
    pythoncom.CoInitialize()
    word = doc = None
    try:
        word = win32com.client.Dispatch("Word.Application"); word.Visible = False
        doc  = word.Documents.Open(os.path.abspath(entrada))
        doc.SaveAs(os.path.abspath(saida), FileFormat=17)
    finally:
        if doc:  
            try: doc.Close(False)
            except Exception: pass
        if word: 
            try: word.Quit()
            except Exception: pass
        pythoncom.CoUninitialize()


# ─── PPT/PPTX → PDF ──────────────────────────────────────────

def ppt_para_pdf(entrada, saida, orientacao="retrato"):
    if MOTOR == "office":        _office_ppt_para_pdf(entrada, saida)
    elif MOTOR == "libreoffice": _soffice_convert(entrada, saida, "pdf")
    else: raise RuntimeError("Nenhum motor disponível.")

def _office_ppt_para_pdf(entrada, saida):
    pythoncom.CoInitialize()
    pp = apres = None
    try:
        pp    = win32com.client.Dispatch("PowerPoint.Application"); pp.Visible = 1
        apres = pp.Presentations.Open(os.path.abspath(entrada), WithWindow=False)
        apres.SaveAs(os.path.abspath(saida), 32)
    finally:
        if apres: 
            try: apres.Close()
            except Exception: pass
        if pp:    
            try: pp.Quit()
            except Exception: pass
        pythoncom.CoUninitialize()


# ─── PDF → DOCX ──────────────────────────────────────────────────

def _pdf2docx_worker(entrada: str, saida: str, inicio: int = 0, fim: int = None):
    """
    Worker isolado para pdf2docx — executado em subprocess separado para
    garantir que pode ser morto por timeout sem travar o processo principal.
    """
    import sys
    try:
        # Desabilita threads do OpenCV para reduzir RAM
        try:
            import cv2
            cv2.setNumThreads(0)
        except Exception:
            pass

        from pdf2docx import Converter as PDF2DOCXConverter
        cv = PDF2DOCXConverter(entrada)
        try:
            kwargs = {"start": inicio}
            if fim is not None:
                kwargs["end"] = fim
            cv.convert(saida, **kwargs)
        finally:
            cv.close()

        if not (os.path.exists(saida) and os.path.getsize(saida) > 0):
            sys.exit(2)  # arquivo vazio = erro
        sys.exit(0)
    except Exception as ex:
        print(f"pdf2docx error: {ex}", file=sys.stderr)
        sys.exit(1)


def pdf_para_docx(entrada, saida):
    """
    Converte PDF para DOCX usando estratégia multicamadas resiliente:
    1. pdf2docx em subprocess com timeout de 90s (isolado para poder matar)
    2. LibreOffice Headless (fallback nativo do sistema)
    3. PyMuPDF + python-docx (fallback de emergência — texto puro)

    Para PDFs grandes (>10 páginas), a conversão com pdf2docx é feita em
    blocos de 10 páginas para evitar estouro de memória no Cloud Run.
    """
    import multiprocessing
    import tempfile

    PDF2DOCX_TIMEOUT = 90  # segundos por tentativa
    MAX_PAGINAS_BLOCO = 10  # páginas por bloco para PDFs grandes

    # Detectar número de páginas do PDF
    try:
        doc_info = fitz.open(entrada)
        total_paginas = doc_info.page_count
        doc_info.close()
    except Exception:
        total_paginas = 0

    log.info("[pdf_para_docx] PDF com %d página(s): %s", total_paginas, os.path.basename(entrada))

    # 1ª Tentativa: pdf2docx via subprocess com timeout
    try:
        if total_paginas <= MAX_PAGINAS_BLOCO:
            # Conversão direta para PDFs pequenos
            p = multiprocessing.Process(
                target=_pdf2docx_worker,
                args=(os.path.abspath(entrada), os.path.abspath(saida))
            )
            p.start()
            p.join(timeout=PDF2DOCX_TIMEOUT)

            if p.is_alive():
                p.terminate()
                p.join(timeout=5)
                if p.is_alive():
                    p.kill()
                log.warning("[pdf_para_docx] pdf2docx ultrapassou timeout de %ds. Usando fallback.", PDF2DOCX_TIMEOUT)
            elif p.exitcode == 0 and os.path.exists(saida) and os.path.getsize(saida) > 0:
                log.info("[pdf_para_docx] Conversão concluída via pdf2docx.")
                return
            else:
                log.warning("[pdf_para_docx] pdf2docx falhou (exit=%s). Iniciando fallback.", p.exitcode)

        else:
            # PDFs grandes: divide em blocos e une no final
            log.info("[pdf_para_docx] PDF grande (%d págs). Convertendo em blocos de %d.", total_paginas, MAX_PAGINAS_BLOCO)
            from docx import Document
            from docx.shared import Pt
            import copy

            doc_final = Document()
            doc_dir = os.path.dirname(os.path.abspath(saida))
            blocos_ok = 0

            for inicio in range(0, total_paginas, MAX_PAGINAS_BLOCO):
                fim = min(inicio + MAX_PAGINAS_BLOCO, total_paginas)
                temp_saida = os.path.join(doc_dir, f"_bloco_{inicio}_{fim}.docx")

                p = multiprocessing.Process(
                    target=_pdf2docx_worker,
                    args=(os.path.abspath(entrada), temp_saida, inicio, fim)
                )
                p.start()
                p.join(timeout=PDF2DOCX_TIMEOUT)

                if p.is_alive():
                    p.terminate()
                    p.join(timeout=5)
                    if p.is_alive(): p.kill()
                    log.warning("[pdf_para_docx] Bloco %d-%d: timeout. Pulando.", inicio, fim)
                    continue

                if p.exitcode == 0 and os.path.exists(temp_saida) and os.path.getsize(temp_saida) > 0:
                    # Mesclar bloco no documento final
                    try:
                        bloco = Document(temp_saida)
                        if blocos_ok > 0:
                            doc_final.add_page_break()
                        for elemento in bloco.paragraphs:
                            p_novo = doc_final.add_paragraph()
                            for run in elemento.runs:
                                r_novo = p_novo.add_run(run.text)
                                r_novo.bold = run.bold
                                r_novo.italic = run.italic
                        blocos_ok += 1
                    except Exception as ex_merge:
                        log.warning("[pdf_para_docx] Erro ao mesclar bloco %d-%d: %s", inicio, fim, ex_merge)
                    finally:
                        try: os.remove(temp_saida)
                        except Exception: pass

            if blocos_ok > 0:
                doc_final.save(saida)
                if os.path.exists(saida) and os.path.getsize(saida) > 0:
                    log.info("[pdf_para_docx] Conversão em blocos concluída (%d blocos).", blocos_ok)
                    return
    except Exception as e:
        log.warning("[pdf_para_docx] Falha na etapa pdf2docx: %s. Iniciando fallback.", e)

    # 2ª Tentativa (Fallback): LibreOffice
    if SOFFICE:
        try:
            log.info("[pdf_para_docx] Tentando conversão via LibreOffice...")
            _soffice_convert(entrada, saida, "docx")
            if os.path.exists(saida) and os.path.getsize(saida) > 0:
                log.info("[pdf_para_docx] Conversão concluída via LibreOffice.")
                return
            log.warning("[pdf_para_docx] LibreOffice não gerou arquivo válido.")
        except Exception as e_soffice:
            log.warning("[pdf_para_docx] Fallback LibreOffice falhou: %s", e_soffice)

    # 3ª Tentativa (Emergência): PyMuPDF + python-docx (texto puro)
    try:
        log.info("[pdf_para_docx] Tentando fallback de emergência (extração de texto)...")
        from docx import Document
        from docx.shared import Pt

        doc_pdf = fitz.open(entrada)
        doc_docx = Document()

        # Estilo de parágrafo simples para melhor legibilidade
        estilo = doc_docx.styles["Normal"]
        estilo.font.size = Pt(11)
        estilo.font.name = "Calibri"

        texto_encontrado = False
        for i, page in enumerate(doc_pdf):
            blocos = page.get_text("blocks")  # Retorna lista de blocos com coordenadas
            if blocos:
                if i > 0:
                    doc_docx.add_page_break()
                for bloco in blocos:
                    texto_bloco = bloco[4].strip()  # bloco[4] = texto do bloco
                    if texto_bloco:
                        doc_docx.add_paragraph(texto_bloco)
                        texto_encontrado = True

        doc_pdf.close()

        if texto_encontrado:
            doc_docx.save(saida)
            if os.path.exists(saida) and os.path.getsize(saida) > 0:
                log.info("[pdf_para_docx] Conversão concluída via fallback PyMuPDF (texto).")
                return
        else:
            log.warning("[pdf_para_docx] PDF sem texto selecionável (provavelmente scanneado).")
            raise RuntimeError(
                "Este PDF parece ser um documento digitalizado (imagem). "
                "Não é possível extrair texto. Tente um PDF com texto selecionável."
            )
    except RuntimeError:
        raise
    except Exception as e_emergencia:
        log.error("[pdf_para_docx] Fallback de emergência falhou: %s", e_emergencia)

    raise RuntimeError(
        "Não foi possível converter o arquivo PDF para Word (DOCX). "
        "Verifique se o PDF contém texto selecionável e não está corrompido."
    )



# ─── PDF → PNG (primeira página) ─────────────────────────────

def pdf_para_png(entrada, saida):
    doc = fitz.open(entrada)
    try:
        pix = doc[0].get_pixmap(matrix=fitz.Matrix(2, 2))
        pix.save(saida)
    finally: doc.close()


# ─── PDF → JPG ────────────────────────────────────────────────

def pdf_para_jpg(entrada, saida):
    doc = fitz.open(entrada)
    try:
        pix = doc[0].get_pixmap(matrix=fitz.Matrix(2, 2))
        img = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")
        img.save(saida, "JPEG", quality=92)
    finally: doc.close()


# ─── PDF → PPTX (cada página = 1 slide) ──────────────────────

def pdf_para_pptx(entrada: str, saida: str):
    from pptx import Presentation
    from pptx.util import Emu
    import tempfile

    doc = fitz.open(entrada)
    prs = Presentation()

    primeira  = doc[0].rect
    ratio     = (primeira.height / primeira.width) if primeira.width else 1.0
    LARG_EMU  = 9144000
    prs.slide_width  = Emu(LARG_EMU)
    prs.slide_height = Emu(int(LARG_EMU * ratio))
    layout_branco = prs.slide_layouts[6]
    temp_dir = tempfile.mkdtemp()

    try:
        for i, page in enumerate(doc):
            mat = fitz.Matrix(2.5, 2.5)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img = os.path.join(temp_dir, f"pg_{i:04d}.png")
            pix.save(img)
            slide = prs.slides.add_slide(layout_branco)
            slide.shapes.add_picture(img, 0, 0,
                                     width=prs.slide_width,
                                     height=prs.slide_height)
    finally:
        doc.close()
        for f in os.listdir(temp_dir):
            try: os.remove(os.path.join(temp_dir, f))
            except Exception: pass
        try: os.rmdir(temp_dir)
        except Exception: pass

    prs.save(saida)


# ─── PDF → PPT (via PPTX + motor) ────────────────────────────

def pdf_para_ppt(entrada: str, saida: str):
    saida_dir = os.path.dirname(saida) or "."
    temp_pptx = os.path.join(saida_dir, f"_tmp_{uuid.uuid4().hex}.pptx")
    try:
        pdf_para_pptx(entrada, temp_pptx)
        if MOTOR == "office":
            _office_pptx_para_ppt(temp_pptx, saida)
        elif MOTOR == "libreoffice":
            _soffice_convert(temp_pptx, saida, "ppt")
        else:
            raise RuntimeError("Nenhum motor disponível para gerar PPT.")
    finally:
        if os.path.exists(temp_pptx):
            try: os.remove(temp_pptx)
            except Exception: pass

def _office_pptx_para_ppt(entrada, saida):
    pythoncom.CoInitialize()
    pp = apres = None
    try:
        pp    = win32com.client.Dispatch("PowerPoint.Application"); pp.Visible = 1
        apres = pp.Presentations.Open(os.path.abspath(entrada), WithWindow=False)
        apres.SaveAs(os.path.abspath(saida), 1)   # 1 = PPT 97-2003
    finally:
        if apres: 
            try: apres.Close()
            except Exception: pass
        if pp:    
            try: pp.Quit()
            except Exception: pass
        pythoncom.CoUninitialize()


# ─── PDF → tabela (3 estratégias) ────────────────────────────

def pdf_extrair_dataframe(entrada: str):
    try:
        import pdfplumber
        import pandas as pd
    except ImportError:
        raise RuntimeError("pdfplumber ou pandas não instalados.")

    # 1. Bordas visíveis
    dados = _tentar_extrair(entrada, "lines")
    if dados:
        df = _normalizar(dados)
        if len(df.columns) > 1: return df

    # 2. Texto agrupado
    dados = _tentar_extrair(entrada, "text")
    if dados:
        df = _normalizar(dados)
        if len(df.columns) > 1: return df

    # 3. Posição X das palavras
    df = _extrair_por_posicao_palavras(entrada)
    if df is not None and len(df.columns) > 1: return df

    # 4. Fallback texto
    return _fallback_texto(entrada)


def _tentar_extrair(entrada: str, estrategia: str) -> list:
    # AUD-002: import explícito (antes dependia do escopo local de pdf_extrair_dataframe)
    import pdfplumber
    cfg = (
        dict(vertical_strategy="lines", horizontal_strategy="lines",
             snap_tolerance=5, join_tolerance=3)
        if estrategia == "lines" else
        dict(vertical_strategy="text",  horizontal_strategy="text",
             snap_tolerance=5, join_tolerance=5,
             intersection_x_tolerance=15, intersection_y_tolerance=15)
    )
    todas = []
    with pdfplumber.open(entrada) as pdf:
        for page in pdf.pages:
            try:
                for t in page.extract_tables(cfg):
                    if t: todas.extend(t)
            except Exception: pass
    return todas


def _normalizar(linhas: list):
    try: import pandas as pd
    except ImportError: raise RuntimeError("pandas não instalado")
    linhas = [r for r in linhas if any(c for c in r if c and str(c).strip())]
    if not linhas: return pd.DataFrame({"Conteúdo": []})
    max_cols = max(len(r) for r in linhas)
    header   = [(str(c).strip() if c else f"Col_{i}") for i, c in enumerate(linhas[0])]
    rows     = [(r + [None]*max_cols)[:max_cols] for r in linhas[1:]]
    df = pd.DataFrame(rows, columns=header)
    df = df.fillna("").map(lambda x: str(x).strip())
    df = df.loc[:, (df != "").any(axis=0)]
    return df


def _extrair_por_posicao_palavras(entrada: str):
    """Detecta colunas pela posição X das palavras — funciona para relatórios sem bordas."""
    import pdfplumber  # AUD-002
    try: import pandas as pd
    except ImportError: raise RuntimeError("pandas não instalado")
    todas_linhas = []

    with pdfplumber.open(entrada) as pdf:
        for page in pdf.pages:
            palavras = page.extract_words(x_tolerance=3, y_tolerance=5, keep_blank_chars=False)
            if not palavras: continue
            linha_atual = [palavras[0]]
            for p in palavras[1:]:
                if abs(p['top'] - linha_atual[0]['top']) < 6:
                    linha_atual.append(p)
                else:
                    todas_linhas.append(sorted(linha_atual, key=lambda w: w['x0']))
                    linha_atual = [p]
            if linha_atual:
                todas_linhas.append(sorted(linha_atual, key=lambda w: w['x0']))

    if not todas_linhas: return None

    contagem_gaps = {}
    for linha in todas_linhas:
        for i in range(len(linha) - 1):
            gap = linha[i+1]['x0'] - linha[i]['x1']
            if gap > 10:
                pos = round((linha[i]['x1'] + linha[i+1]['x0']) / 2 / 15) * 15
                contagem_gaps[pos] = contagem_gaps.get(pos, 0) + 1

    if not contagem_gaps: return None

    min_freq     = max(2, len(todas_linhas) * 0.15)
    separadores  = sorted(pos for pos, cnt in contagem_gaps.items() if cnt >= min_freq)
    if not separadores: return None

    seps = [separadores[0]]
    for s in separadores[1:]:
        if s - seps[-1] > 35: seps.append(s)

    limites = [0] + seps + [float('inf')]
    n_cols  = len(limites) - 1
    if n_cols < 2: return None

    matrix = []
    for linha in todas_linhas:
        row = [""] * n_cols
        for p in linha:
            x_mid = (p['x0'] + p['x1']) / 2
            for i in range(n_cols):
                if limites[i] <= x_mid < limites[i+1]:
                    row[i] = (row[i] + " " + p['text']).strip()
                    break
        if any(c for c in row): matrix.append(row)

    if len(matrix) < 2: return None

    header = [(matrix[0][i] if matrix[0][i].strip() else f"Col_{i}") for i in range(n_cols)]
    df = pd.DataFrame(matrix[1:], columns=header)
    df = df.fillna("").map(lambda x: str(x).strip())
    df = df.loc[:, (df != "").any(axis=0)]
    return df


def _fallback_texto(entrada: str):
    import pdfplumber  # AUD-002
    try: import pandas as pd
    except ImportError: raise RuntimeError("pandas não instalado")
    linhas_raw = []
    with pdfplumber.open(entrada) as pdf:
        for page in pdf.pages:
            txt = page.extract_text(x_tolerance=3, y_tolerance=3)
            if txt:
                for l in txt.split("\n"):
                    if l.strip(): linhas_raw.append(l.strip())
    if not linhas_raw:
        raise ValueError("Nenhum dado extraível. O PDF pode ser uma imagem escaneada.")
    max_cols = 1
    split    = []
    for l in linhas_raw:
        parts = [p.strip() for p in _re.split(r" {2,}", l) if p.strip()]
        split.append(parts)
        if len(parts) > max_cols: max_cols = len(parts)
    if max_cols == 1:
        return pd.DataFrame(linhas_raw, columns=["Conteúdo"])
    header = [(split[0][i] if i < len(split[0]) else f"Col_{i}") for i in range(max_cols)]
    rows   = [(r + [""]*max_cols)[:max_cols] for r in split[1:]]
    return pd.DataFrame(rows, columns=header)


def pdf_para_xlsx(entrada: str, saida: str):
    df = pdf_extrair_dataframe(entrada)
    df.to_excel(saida, index=False, engine="openpyxl")

def pdf_para_csv(entrada: str, saida: str):
    df = pdf_extrair_dataframe(entrada)
    df.to_csv(saida, index=False, encoding="utf-8-sig")


# ─── Imagem → PDF ────────────────────────────────────────────

def imagem_para_pdf(entrada, saida):
    img = Image.open(entrada)
    if img.mode in ("RGBA", "LA", "P"): img = img.convert("RGB")
    img.save(saida, "PDF", resolution=150)


# ─── PNG ↔ JPG ────────────────────────────────────────────────

def png_para_jpg(entrada, saida):
    img = Image.open(entrada).convert("RGB")
    img.save(saida, "JPEG", quality=92)

def jpg_para_png(entrada, saida):
    img = Image.open(entrada).convert("RGBA")
    img.save(saida, "PNG")


# ─── DOCX/XLSX → Imagem (via PDF) ────────────────────────────

def _pdf_para_imagem_completo(pdf_path: str, saida: str, fmt: str = "png"):
    """Renderiza todas as páginas e empacota em um arquivo .zip para evitar Out-of-Memory (OOM)."""
    # AUD-008: io e zipfile já importados no topo do módulo
    doc = fitz.open(pdf_path)
    try:
        with zipfile.ZipFile(saida, 'w', zipfile.ZIP_DEFLATED) as zipf:
            for i, page in enumerate(doc):
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img_name = f"pagina_{i+1:03d}.{fmt}"
                
                if fmt in ("jpg", "webp", "heic"):
                    img = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")
                    img_io = io.BytesIO()
                    img.save(img_io, "JPEG" if fmt == "jpg" else fmt.upper(), quality=92)
                    zipf.writestr(img_name, img_io.getvalue())
                else:
                    zipf.writestr(img_name, pix.tobytes("png"))
    finally:
        doc.close()

def _via_pdf_para_img(entrada: str, saida: str, origem_ext: str, fmt: str):
    saida_dir = os.path.dirname(saida) or "."
    temp_pdf  = os.path.join(saida_dir, f"_tmp_{uuid.uuid4().hex}.pdf")
    try:
        if origem_ext in ("docx", "doc"):   docx_para_pdf(entrada, temp_pdf)
        elif origem_ext in ("xlsx", "xls"): xlsx_para_pdf(entrada, temp_pdf)
        elif origem_ext in ("csv",):        csv_para_pdf(entrada, temp_pdf)
        elif origem_ext in ("json",):       json_para_pdf(entrada, temp_pdf)
        elif origem_ext in ("ppt",):        ppt_para_pdf(entrada, temp_pdf)
        elif origem_ext in ("pptx",):       ppt_para_pdf(entrada, temp_pdf)
        else: raise ValueError(f"Origem '{origem_ext}' não suportada para imagem.")
        _pdf_para_imagem_completo(temp_pdf, saida, fmt=fmt)
    finally:
        if os.path.exists(temp_pdf):
            try: os.remove(temp_pdf)
            except Exception: pass

def docx_para_png(entrada, saida): _via_pdf_para_img(entrada, saida, "docx", "png")
def docx_para_jpg(entrada, saida): _via_pdf_para_img(entrada, saida, "docx", "jpg")
def docx_para_webp(entrada, saida): _via_pdf_para_img(entrada, saida, "docx", "webp")
def docx_para_heic(entrada, saida): _via_pdf_para_img(entrada, saida, "docx", "heic")

def xlsx_para_png(entrada, saida): _via_pdf_para_img(entrada, saida, "xlsx", "png")
def xlsx_para_jpg(entrada, saida): _via_pdf_para_img(entrada, saida, "xlsx", "jpg")
def xlsx_para_webp(entrada, saida): _via_pdf_para_img(entrada, saida, "xlsx", "webp")
def xlsx_para_heic(entrada, saida): _via_pdf_para_img(entrada, saida, "xlsx", "heic")

def csv_para_png(entrada, saida):  _via_pdf_para_img(entrada, saida, "csv", "png")
def csv_para_jpg(entrada, saida):  _via_pdf_para_img(entrada, saida, "csv", "jpg")
def csv_para_webp(entrada, saida): _via_pdf_para_img(entrada, saida, "csv", "webp")
def csv_para_heic(entrada, saida): _via_pdf_para_img(entrada, saida, "csv", "heic")

def json_para_png(entrada, saida):  _via_pdf_para_img(entrada, saida, "json", "png")
def json_para_jpg(entrada, saida):  _via_pdf_para_img(entrada, saida, "json", "jpg")
def json_para_webp(entrada, saida): _via_pdf_para_img(entrada, saida, "json", "webp")
def json_para_heic(entrada, saida): _via_pdf_para_img(entrada, saida, "json", "heic")

def ppt_para_png(entrada, saida):  _via_pdf_para_img(entrada, saida, "ppt", "png")
def ppt_para_jpg(entrada, saida):  _via_pdf_para_img(entrada, saida, "ppt", "jpg")
def ppt_para_webp(entrada, saida): _via_pdf_para_img(entrada, saida, "ppt", "webp")
def ppt_para_heic(entrada, saida): _via_pdf_para_img(entrada, saida, "ppt", "heic")

def pptx_para_png(entrada, saida): _via_pdf_para_img(entrada, saida, "pptx", "png")
def pptx_para_jpg(entrada, saida): _via_pdf_para_img(entrada, saida, "pptx", "jpg")
def pptx_para_webp(entrada, saida): _via_pdf_para_img(entrada, saida, "pptx", "webp")
def pptx_para_heic(entrada, saida): _via_pdf_para_img(entrada, saida, "pptx", "heic")

def pdf_para_png(entrada, saida):  _pdf_para_imagem_completo(entrada, saida, fmt="png")
def pdf_para_jpg(entrada, saida):  _pdf_para_imagem_completo(entrada, saida, fmt="jpg")
def pdf_para_webp(entrada, saida): _pdf_para_imagem_completo(entrada, saida, fmt="webp")
def pdf_para_heic(entrada, saida): _pdf_para_imagem_completo(entrada, saida, fmt="heic")


# ─── Conversão genérica via PDF intermediário ─────────────────

def _via_pdf(entrada, saida, origem_ext, destino_ext):
    """Converte qualquer formato para qualquer outro usando PDF como intermediário."""
    saida_dir = os.path.dirname(saida) or "."
    temp_pdf  = os.path.join(saida_dir, f"_tmp_{uuid.uuid4().hex}.pdf")
    try:
        # Passo 1: origem → PDF
        funcao_para_pdf = _MAPA_PARA_PDF.get(origem_ext)
        if not funcao_para_pdf:
            raise ValueError(f"Não é possível converter '{origem_ext}' para PDF.")
        funcao_para_pdf(entrada, temp_pdf)

        # Passo 2: PDF → destino
        funcao_de_pdf = _MAPA_DE_PDF.get(destino_ext)
        if not funcao_de_pdf:
            raise ValueError(f"Não é possível converter PDF para '{destino_ext}'.")
        funcao_de_pdf(temp_pdf, saida)
    finally:
        if os.path.exists(temp_pdf):
            try: os.remove(temp_pdf)
            except Exception: pass


# ─── JSON ↔ CSV/XLSX ──────────────────────────────────────────────
def json_para_csv(entrada, saida):
    try: import pandas as pd
    except ImportError: raise RuntimeError("pandas não instalado")
    df = pd.read_json(entrada)
    df.fillna("").to_csv(saida, index=False, encoding="utf-8-sig")

def csv_para_json(entrada, saida):
    try: import pandas as pd
    except ImportError: raise RuntimeError("pandas não instalado")
    enc = detectar_encoding(entrada)
    df = _carregar_csv(entrada, enc)
    df.fillna("").to_json(saida, orient="records", force_ascii=False, indent=2)

def json_para_xlsx(entrada, saida):
    try: import pandas as pd
    except ImportError: raise RuntimeError("pandas não instalado")
    df = pd.read_json(entrada)
    df.fillna("").to_excel(saida, index=False, engine="openpyxl")

def xlsx_para_json(entrada, saida):
    try: import pandas as pd
    except ImportError: raise RuntimeError("pandas não instalado")
    df = pd.read_excel(entrada, engine="openpyxl")
    df.fillna("").to_json(saida, orient="records", force_ascii=False, indent=2)

def json_para_pdf(entrada, saida, orientacao="retrato"):
    temp = os.path.join(os.path.dirname(saida) or ".", f"_tmp_{uuid.uuid4().hex}.xlsx")
    try:
        json_para_xlsx(entrada, temp)
        xlsx_para_pdf(temp, saida, orientacao=orientacao)
    finally:
        if os.path.exists(temp):
            try: os.remove(temp)
            except Exception: pass





# ─── Mesclar Planilhas ──────────────────────────────────────────
def mesclar_planilhas(arquivos: list, saida: str, formato: str = "xlsx"):
    try: import pandas as pd
    except ImportError: raise RuntimeError("pandas não instalado")
    dfs = []
    for arquivo in arquivos:
        ext = arquivo.split(".")[-1].lower()
        if ext == "csv":
            dfs.append(_carregar_csv(arquivo, detectar_encoding(arquivo)))
        elif ext in ("xlsx", "xls"):
            dfs.append(pd.read_excel(arquivo, engine="openpyxl"))
        elif ext == "json":
            dfs.append(pd.read_json(arquivo))
    
    if not dfs:
        raise ValueError("Nenhum dado válido para mesclar.")
        
    df_final = pd.concat(dfs, ignore_index=True)
    
    if formato == "csv":
        df_final.to_csv(saida, index=False, encoding="utf-8-sig")
    elif formato == "json":
        df_final.to_json(saida, orient="records", force_ascii=False, indent=2)
    else:
        df_final.to_excel(saida, index=False, engine="openpyxl")

# ─── OCR ───────────────────────────────────────────────
def imagem_para_txt_ocr(entrada: str, saida: str):
    # pyrefly: ignore [missing-import]
    try: import pytesseract
    except ImportError: raise RuntimeError("pytesseract não instalado.")
    img = Image.open(entrada).convert("RGB")
    texto = pytesseract.image_to_string(img, lang="por+eng")
    with open(saida, "w", encoding="utf-8") as f:
        f.write(texto)

def pdf_para_txt_ocr(entrada: str, saida: str):
    # pyrefly: ignore [missing-import]
    try: import pytesseract
    except ImportError: raise RuntimeError("pytesseract não instalado.")
    doc = fitz.open(entrada)
    texto_total = []
    try:
        for page in doc:
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            img = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")
            texto = pytesseract.image_to_string(img, lang="por+eng")
            texto_total.append(texto)
    finally:
        doc.close()
    with open(saida, "w", encoding="utf-8") as f:
        f.write("\n\n---\n\n".join(texto_total))

# ─── WEBP / HEIC ────────────────────────────────────────────────
def imagem_para_webp(entrada, saida):
    img = Image.open(entrada).convert("RGBA")
    img.save(saida, "WEBP", quality=92)

def imagem_para_heic(entrada, saida):
    img = Image.open(entrada).convert("RGB")
    img.save(saida, "HEIC", quality=92)

def heic_para_png(entrada, saida):
    img = Image.open(entrada).convert("RGBA")
    img.save(saida, "PNG")

def heic_para_jpg(entrada, saida):
    img = Image.open(entrada).convert("RGB")
    img.save(saida, "JPEG", quality=92)


# Mapas auxiliares para conversão via PDF
_MAPA_PARA_PDF = {
    "csv":  lambda e, s: csv_para_pdf(e, s),
    "xlsx": lambda e, s: xlsx_para_pdf(e, s),
    "xls":  lambda e, s: xlsx_para_pdf(e, s),
    "json": lambda e, s: json_para_pdf(e, s),
    "docx": lambda e, s: docx_para_pdf(e, s),
    "doc":  lambda e, s: docx_para_pdf(e, s),
    "ppt":  lambda e, s: ppt_para_pdf(e, s),
    "pptx": lambda e, s: ppt_para_pdf(e, s),
    "png":  imagem_para_pdf,
    "jpg":  imagem_para_pdf,
    "jpeg": imagem_para_pdf,
    "webp": imagem_para_pdf,
    "heic": imagem_para_pdf,
}

_MAPA_DE_PDF = {
    "docx": pdf_para_docx,
    "pptx": pdf_para_pptx,
    "ppt":  pdf_para_ppt,
    "xlsx": pdf_para_xlsx,
    "csv":  pdf_para_csv,
    "png":  pdf_para_png,
    "jpg":  pdf_para_jpg,
    "webp": pdf_para_webp,
    "heic": pdf_para_heic,
}

def csv_para_docx(entrada, saida):
    try:
        import pandas as pd
        from docx import Document
    except ImportError:
        _via_pdf(entrada, saida, "csv", "docx")
        return
    enc = detectar_encoding(entrada)
    df = _carregar_csv(entrada, enc).fillna("")
    doc = Document()
    t = doc.add_table(rows=len(df) + 1, cols=len(df.columns))
    t.style = 'Table Grid'
    for j, col_name in enumerate(df.columns):
        t.cell(0, j).text = str(col_name)
    for i, row in enumerate(df.itertuples(index=False)):
        for j, val in enumerate(row):
            t.cell(i + 1, j).text = str(val)
    doc.save(saida)

def xlsx_para_docx(entrada, saida):
    try:
        import pandas as pd
        from docx import Document
    except ImportError:
        _via_pdf(entrada, saida, "xlsx", "docx")
        return
    df = pd.read_excel(entrada, engine="openpyxl").fillna("")
    doc = Document()
    t = doc.add_table(rows=len(df) + 1, cols=len(df.columns))
    t.style = 'Table Grid'
    for j, col_name in enumerate(df.columns):
        t.cell(0, j).text = str(col_name)
    for i, row in enumerate(df.itertuples(index=False)):
        for j, val in enumerate(row):
            t.cell(i + 1, j).text = str(val)
    doc.save(saida)

def csv_para_pptx(e, s):  _via_pdf(e, s, "csv",  "pptx")
def xlsx_para_pptx(e, s): _via_pdf(e, s, "xlsx", "pptx")
def json_para_docx(e, s): _via_pdf(e, s, "json", "docx")
def json_para_pptx(e, s): _via_pdf(e, s, "json", "pptx")
def docx_para_xlsx(e, s): _via_pdf(e, s, "docx", "xlsx")
def docx_para_csv(e, s):  _via_pdf(e, s, "docx", "csv")
def docx_para_pptx(e, s): _via_pdf(e, s, "docx", "pptx")
def ppt_para_docx(e, s):  _via_pdf(e, s, "ppt",  "docx")
def ppt_para_xlsx(e, s):  _via_pdf(e, s, "ppt",  "xlsx")
def ppt_para_csv(e, s):   _via_pdf(e, s, "ppt",  "csv")
def ppt_para_pptx(e, s):  _via_pdf(e, s, "ppt",  "pptx")
def pptx_para_docx(e, s): _via_pdf(e, s, "pptx", "docx")
def pptx_para_xlsx(e, s): _via_pdf(e, s, "pptx", "xlsx")
def pptx_para_csv(e, s):  _via_pdf(e, s, "pptx", "csv")
def pptx_para_ppt(e, s):  _via_pdf(e, s, "pptx", "ppt")
def png_para_docx(e, s):  _via_pdf(e, s, "png",  "docx")
def png_para_pptx(e, s):  _via_pdf(e, s, "png",  "pptx")
def jpg_para_docx(e, s):  _via_pdf(e, s, "jpg",  "docx")
def jpg_para_pptx(e, s):  _via_pdf(e, s, "jpg",  "pptx")
def webp_para_docx(e, s): _via_pdf(e, s, "webp", "docx")
def webp_para_pptx(e, s): _via_pdf(e, s, "webp", "pptx")
def heic_para_docx(e, s): _via_pdf(e, s, "heic", "docx")
def heic_para_pptx(e, s): _via_pdf(e, s, "heic", "pptx")
def jpg_para_png_img(e, s): jpg_para_png(e, s)


# ─── Despacho central ─────────────────────────────────────────

_MAPA = {
    # CSV
    ("csv",  "xlsx"): csv_para_xlsx,
    ("csv",  "pdf"):  csv_para_pdf,
    ("csv",  "png"):  csv_para_png,
    ("csv",  "jpg"):  csv_para_jpg,
    ("csv",  "webp"): csv_para_webp,
    ("csv",  "heic"): csv_para_heic,
    ("csv",  "docx"): csv_para_docx,
    ("csv",  "pptx"): csv_para_pptx,
    ("csv",  "json"): csv_para_json,
    # XLSX
    ("xlsx", "csv"):  xlsx_para_csv,
    ("xlsx", "pdf"):  xlsx_para_pdf,
    ("xlsx", "png"):  xlsx_para_png,
    ("xlsx", "jpg"):  xlsx_para_jpg,
    ("xlsx", "webp"): xlsx_para_webp,
    ("xlsx", "heic"): xlsx_para_heic,
    ("xlsx", "docx"): xlsx_para_docx,
    ("xlsx", "pptx"): xlsx_para_pptx,
    ("xlsx", "json"): xlsx_para_json,
    # JSON
    ("json", "csv"):  json_para_csv,
    ("json", "xlsx"): json_para_xlsx,
    ("json", "pdf"):  json_para_pdf,
    ("json", "png"):  json_para_png,
    ("json", "jpg"):  json_para_jpg,
    ("json", "webp"): json_para_webp,
    ("json", "heic"): json_para_heic,
    # PDF
    ("pdf",  "docx"): pdf_para_docx,
    ("pdf",  "pptx"): pdf_para_pptx,
    ("pdf",  "ppt"):  pdf_para_ppt,
    ("pdf",  "png"):  pdf_para_png,
    ("pdf",  "jpg"):  pdf_para_jpg,
    ("pdf",  "webp"): pdf_para_webp,
    ("pdf",  "heic"): pdf_para_heic,
    ("pdf",  "xlsx"): pdf_para_xlsx,
    ("pdf",  "csv"):  pdf_para_csv,
    ("pdf",  "txt"):  pdf_para_txt_ocr,
    # DOCX
    ("docx", "pdf"):  docx_para_pdf,
    ("docx", "png"):  docx_para_png,
    ("docx", "jpg"):  docx_para_jpg,
    ("docx", "webp"): docx_para_webp,
    ("docx", "heic"): docx_para_heic,
    ("docx", "xlsx"): docx_para_xlsx,
    ("docx", "csv"):  docx_para_csv,
    ("docx", "pptx"): docx_para_pptx,
    # PPT
    ("ppt",  "pdf"):  ppt_para_pdf,
    ("ppt",  "docx"): ppt_para_docx,
    ("ppt",  "xlsx"): ppt_para_xlsx,
    ("ppt",  "csv"):  ppt_para_csv,
    ("ppt",  "png"):  ppt_para_png,
    ("ppt",  "jpg"):  ppt_para_jpg,
    ("ppt",  "webp"): ppt_para_webp,
    ("ppt",  "heic"): ppt_para_heic,
    ("ppt",  "pptx"): ppt_para_pptx,
    # PPTX
    ("pptx", "pdf"):  ppt_para_pdf,
    ("pptx", "docx"): pptx_para_docx,
    ("pptx", "xlsx"): pptx_para_xlsx,
    ("pptx", "csv"):  pptx_para_csv,
    ("pptx", "png"):  pptx_para_png,
    ("pptx", "jpg"):  pptx_para_jpg,
    ("pptx", "webp"): pptx_para_webp,
    ("pptx", "heic"): pptx_para_heic,
    ("pptx", "ppt"):  pptx_para_ppt,
    # PNG
    ("png",  "pdf"):  imagem_para_pdf,
    ("png",  "jpg"):  png_para_jpg,
    ("png",  "webp"): imagem_para_webp,
    ("png",  "heic"): imagem_para_heic,
    ("png",  "docx"): png_para_docx,
    ("png",  "pptx"): png_para_pptx,
    ("png",  "txt"):  imagem_para_txt_ocr,
    # JPG
    ("jpg",  "pdf"):  imagem_para_pdf,
    ("jpg",  "png"):  jpg_para_png,
    ("jpg",  "webp"): imagem_para_webp,
    ("jpg",  "heic"): imagem_para_heic,
    ("jpg",  "docx"): jpg_para_docx,
    ("jpg",  "pptx"): jpg_para_pptx,
    ("jpg",  "txt"):  imagem_para_txt_ocr,
    # JPEG (alias)
    ("jpeg", "pdf"):  imagem_para_pdf,
    ("jpeg", "png"):  jpg_para_png,
    ("jpeg", "webp"): imagem_para_webp,
    ("jpeg", "heic"): imagem_para_heic,
    ("jpeg", "docx"): jpg_para_docx,
    ("jpeg", "pptx"): jpg_para_pptx,
    ("jpeg", "txt"):  imagem_para_txt_ocr,
    # WEBP
    ("webp", "pdf"):  imagem_para_pdf,
    ("webp", "png"):  jpg_para_png,
    ("webp", "jpg"):  png_para_jpg,
    ("webp", "heic"): imagem_para_heic,
    ("webp", "docx"): webp_para_docx,
    ("webp", "pptx"): webp_para_pptx,
    ("webp", "txt"):  imagem_para_txt_ocr,
    # HEIC
    ("heic", "pdf"):  imagem_para_pdf,
    ("heic", "png"):  heic_para_png,
    ("heic", "jpg"):  heic_para_jpg,
    ("heic", "webp"): imagem_para_webp,
    ("heic", "docx"): heic_para_docx,
    ("heic", "pptx"): heic_para_pptx,
    ("heic", "txt"):  imagem_para_txt_ocr,

}

_ACEITA_ORIENTACAO = {("xlsx","pdf"), ("xls","pdf"), ("csv","pdf"), ("json","pdf")}

def converter_arquivo(entrada, saida, origem, destino, orientacao="retrato"):
    origem  = origem.lower()
    destino = destino.lower()
    if not os.path.exists(entrada):
        raise FileNotFoundError(f"Arquivo não encontrado: {entrada}")
    funcao = _MAPA.get((origem, destino))
    if not funcao:
        raise ValueError(f"Conversão '{origem.upper()}→{destino.upper()}' não suportada.")
    if (origem, destino) in _ACEITA_ORIENTACAO:
        # pyrefly: ignore [unexpected-keyword]
        funcao(entrada, saida, orientacao=orientacao)
    else:
        funcao(entrada, saida)