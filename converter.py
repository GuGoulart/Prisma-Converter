import os
import re as _re
import uuid
import shutil
import subprocess
import fitz
import pandas as pd
from pdf2docx import Converter
from PIL import Image

try:
    import win32com.client
    import pythoncom
    _TEM_WIN32 = True
except ImportError:
    _TEM_WIN32 = False

try:
    import chardet
    _TEM_CHARDET = True
except ImportError:
    _TEM_CHARDET = False

try:
    import pdfplumber
    _TEM_PDFPLUMBER = True
except ImportError:
    _TEM_PDFPLUMBER = False


# ─── Motor ───────────────────────────────────────────────────

def _encontrar_soffice():
    if shutil.which("soffice"):     return "soffice"
    if shutil.which("libreoffice"): return "libreoffice"
    for c in [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        r"C:\Program Files\LibreOffice 7\program\soffice.exe",
        r"C:\Program Files\LibreOffice 24\program\soffice.exe",
    ]:
        if os.path.exists(c): return c
    return None


SOFFICE = _encontrar_soffice()

if _TEM_WIN32:    MOTOR = "office"
elif SOFFICE:     MOTOR = "libreoffice"
else:             MOTOR = None

print(f"[Prisma] Motor: {MOTOR or 'NENHUM'} | pdfplumber: {_TEM_PDFPLUMBER}")


# ─── Conversões disponíveis ───────────────────────────────────

CONVERSOES = {
    "csv":  ["xlsx", "pdf"],
    "xlsx": ["csv",  "pdf", "png", "jpg"],
    "pdf":  ["docx", "pptx", "ppt", "png", "xlsx", "csv"],
    "docx": ["pdf",  "png", "jpg"],
    "ppt":  ["pdf"],
    "pptx": ["pdf"],
    "png":  ["pdf"],
    "jpg":  ["pdf"],
    "jpeg": ["pdf"],
}

def obter_conversoes(extensao): return CONVERSOES.get(extensao.lower(), [])
def obter_motor():               return MOTOR


# ─── Encoding ────────────────────────────────────────────────

def detectar_encoding(caminho: str) -> str:
    with open(caminho, "rb") as f: raw = f.read(4)
    if raw.startswith(b"\xef\xbb\xbf"):                             return "utf-8-sig"
    if raw.startswith(b"\xff\xfe") or raw.startswith(b"\xfe\xff"): return "utf-16"
    if _TEM_CHARDET:
        with open(caminho, "rb") as f: r = chardet.detect(f.read(50_000))
        return r.get("encoding") or "utf-8"
    for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
        try:
            with open(caminho, encoding=enc) as f: f.read(1000)
            return enc
        except: continue
    return "utf-8"


# ─── LibreOffice headless ─────────────────────────────────────

def _soffice_convert(entrada, saida, fmt, timeout=120):
    if not SOFFICE: raise RuntimeError("LibreOffice não encontrado.")
    saida_dir  = os.path.dirname(os.path.abspath(saida))
    nome_base  = os.path.splitext(os.path.basename(entrada))[0]
    saida_auto = os.path.join(saida_dir, f"{nome_base}.{fmt}")
    r = subprocess.run(
        [SOFFICE, "--headless", "--norestore",
         "--convert-to", fmt, "--outdir", saida_dir,
         os.path.abspath(entrada)],
        capture_output=True, timeout=timeout
    )
    if r.returncode != 0:
        raise RuntimeError(f"LibreOffice: {r.stderr.decode('utf-8', errors='replace')}")
    if os.path.abspath(saida_auto) != os.path.abspath(saida) and os.path.exists(saida_auto):
        os.replace(saida_auto, saida)
    if not os.path.exists(saida):
        raise RuntimeError("LibreOffice não gerou o arquivo de saída.")


# ─── Pré-formata XLSX para PDF organizado ────────────────────

def _preparar_xlsx_para_pdf(caminho: str, orientacao: str = "retrato"):
    from openpyxl import load_workbook
    from openpyxl.worksheet.page import PageMargins
    wb = load_workbook(caminho)
    for ws in wb.worksheets:
        ws.page_setup.orientation = "landscape" if orientacao == "paisagem" else "portrait"
        ws.page_setup.fitToPage   = True
        ws.page_setup.fitToWidth  = 1
        ws.page_setup.fitToHeight = 0
        ws.page_margins = PageMargins(
            left=0.4, right=0.4, top=0.5, bottom=0.5, header=0.2, footer=0.2
        )
        if ws.max_row and ws.max_row > 1:
            ws.print_title_rows = "1:1"
        for col in ws.columns:
            letra   = col[0].column_letter
            max_len = max((len(str(c.value)) if c.value is not None else 0 for c in col), default=0)
            ws.column_dimensions[letra].width = min(max(max_len + 2, 8), 35)
    wb.save(caminho)


# ─── CSV ↔ XLSX ───────────────────────────────────────────────

def csv_para_xlsx(entrada, saida):
    enc = detectar_encoding(entrada)
    df  = pd.read_csv(entrada, encoding=enc, sep=None, engine="python",
                      encoding_errors="replace", on_bad_lines="skip")
    df.fillna("").to_excel(saida, index=False, engine="openpyxl")

def xlsx_para_csv(entrada, saida):
    df = pd.read_excel(entrada, engine="openpyxl")
    df.fillna("").to_csv(saida, index=False, encoding="utf-8-sig")


# ─── XLSX → PDF ───────────────────────────────────────────────

def xlsx_para_pdf(entrada, saida, orientacao="retrato"):
    if MOTOR == "office":
        _office_xlsx_para_pdf(entrada, saida, orientacao)
    elif MOTOR == "libreoffice":
        _preparar_xlsx_para_pdf(entrada, orientacao)
        _soffice_convert(entrada, saida, "pdf")
    else:
        raise RuntimeError("Nenhum motor disponível.")

def _office_xlsx_para_pdf(entrada, saida, orientacao="retrato"):
    pythoncom.CoInitialize()
    excel = wb = None
    try:
        excel = win32com.client.Dispatch("Excel.Application"); excel.Visible = False
        wb    = excel.Workbooks.Open(os.path.abspath(entrada))
        for sheet in wb.Worksheets:
            ps = sheet.PageSetup
            ps.Orientation    = 2 if orientacao == "paisagem" else 1
            ps.Zoom           = False
            ps.FitToPagesWide = 1
            ps.FitToPagesTall = False
            ps.LeftMargin     = excel.InchesToPoints(0.4)
            ps.RightMargin    = excel.InchesToPoints(0.4)
            ps.TopMargin      = excel.InchesToPoints(0.5)
            ps.BottomMargin   = excel.InchesToPoints(0.5)
            ps.PrintTitleRows = "$1:$1"
        wb.ExportAsFixedFormat(0, os.path.abspath(saida))
    finally:
        if wb:    
            try: wb.Close(False)
            except: pass
        if excel: 
            try: excel.Quit()
            except: pass
        pythoncom.CoUninitialize()


# ─── CSV → PDF ────────────────────────────────────────────────

def csv_para_pdf(entrada, saida, orientacao="retrato"):
    enc  = detectar_encoding(entrada)
    temp = os.path.join(os.path.dirname(saida) or ".", f"_tmp_{uuid.uuid4().hex}.xlsx")
    try:
        df = pd.read_csv(entrada, encoding=enc, sep=None, engine="python",
                         encoding_errors="replace", on_bad_lines="skip")
        df.fillna("").to_excel(temp, index=False, engine="openpyxl")
        xlsx_para_pdf(temp, saida, orientacao=orientacao)
    finally:
        if os.path.exists(temp):
            try: os.remove(temp)
            except: pass


# ─── DOCX → PDF ───────────────────────────────────────────────

def docx_para_pdf(entrada, saida, orientacao="retrato"):
    if MOTOR == "office":        _office_docx_para_pdf(entrada, saida)
    elif MOTOR == "libreoffice": _soffice_convert(entrada, saida, "pdf")
    else: raise RuntimeError("Nenhum motor disponível.")

def _office_docx_para_pdf(entrada, saida):
    pythoncom.CoInitialize()
    word = doc = None
    try:
        word = win32com.client.Dispatch("Word.Application"); word.Visible = False
        doc  = word.Documents.Open(os.path.abspath(entrada))
        doc.SaveAs(os.path.abspath(saida), FileFormat=17)
    finally:
        if doc:  
            try: doc.Close(False)
            except: pass
        if word: 
            try: word.Quit()
            except: pass
        pythoncom.CoUninitialize()


# ─── PPT/PPTX → PDF ──────────────────────────────────────────

def ppt_para_pdf(entrada, saida, orientacao="retrato"):
    if MOTOR == "office":        _office_ppt_para_pdf(entrada, saida)
    elif MOTOR == "libreoffice": _soffice_convert(entrada, saida, "pdf")
    else: raise RuntimeError("Nenhum motor disponível.")

def _office_ppt_para_pdf(entrada, saida):
    pythoncom.CoInitialize()
    pp = apres = None
    try:
        pp    = win32com.client.Dispatch("PowerPoint.Application"); pp.Visible = 1
        apres = pp.Presentations.Open(os.path.abspath(entrada), WithWindow=False)
        apres.SaveAs(os.path.abspath(saida), 32)
    finally:
        if apres: 
            try: apres.Close()
            except: pass
        if pp:    
            try: pp.Quit()
            except: pass
        pythoncom.CoUninitialize()


# ─── PDF → DOCX ───────────────────────────────────────────────

def pdf_para_docx(entrada, saida):
    cv = Converter(entrada)
    try:    cv.convert(saida)
    finally: cv.close()


# ─── PDF → PNG (primeira página) ─────────────────────────────

def pdf_para_png(entrada, saida):
    doc = fitz.open(entrada)
    try:
        pix = doc[0].get_pixmap(matrix=fitz.Matrix(2, 2))
        pix.save(saida)
    finally: doc.close()


# ─── PDF → PPTX (cada página = 1 slide) ──────────────────────

def pdf_para_pptx(entrada: str, saida: str):
    from pptx import Presentation
    from pptx.util import Emu
    import tempfile

    doc = fitz.open(entrada)
    prs = Presentation()

    primeira  = doc[0].rect
    ratio     = (primeira.height / primeira.width) if primeira.width else 1.0
    LARG_EMU  = 9144000
    prs.slide_width  = Emu(LARG_EMU)
    prs.slide_height = Emu(int(LARG_EMU * ratio))
    layout_branco = prs.slide_layouts[6]
    temp_dir = tempfile.mkdtemp()

    try:
        for i, page in enumerate(doc):
            mat = fitz.Matrix(2.5, 2.5)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            img = os.path.join(temp_dir, f"pg_{i:04d}.png")
            pix.save(img)
            slide = prs.slides.add_slide(layout_branco)
            slide.shapes.add_picture(img, 0, 0,
                                     width=prs.slide_width,
                                     height=prs.slide_height)
    finally:
        doc.close()
        for f in os.listdir(temp_dir):
            try: os.remove(os.path.join(temp_dir, f))
            except: pass
        try: os.rmdir(temp_dir)
        except: pass

    prs.save(saida)


# ─── PDF → PPT (via PPTX + motor) ────────────────────────────

def pdf_para_ppt(entrada: str, saida: str):
    saida_dir = os.path.dirname(saida) or "."
    temp_pptx = os.path.join(saida_dir, f"_tmp_{uuid.uuid4().hex}.pptx")
    try:
        pdf_para_pptx(entrada, temp_pptx)
        if MOTOR == "office":
            _office_pptx_para_ppt(temp_pptx, saida)
        elif MOTOR == "libreoffice":
            _soffice_convert(temp_pptx, saida, "ppt")
        else:
            raise RuntimeError("Nenhum motor disponível para gerar PPT.")
    finally:
        if os.path.exists(temp_pptx):
            try: os.remove(temp_pptx)
            except: pass

def _office_pptx_para_ppt(entrada, saida):
    pythoncom.CoInitialize()
    pp = apres = None
    try:
        pp    = win32com.client.Dispatch("PowerPoint.Application"); pp.Visible = 1
        apres = pp.Presentations.Open(os.path.abspath(entrada), WithWindow=False)
        apres.SaveAs(os.path.abspath(saida), 1)   # 1 = PPT 97-2003
    finally:
        if apres: 
            try: apres.Close()
            except: pass
        if pp:    
            try: pp.Quit()
            except: pass
        pythoncom.CoUninitialize()


# ─── PDF → tabela (3 estratégias) ────────────────────────────

def pdf_extrair_dataframe(entrada: str) -> pd.DataFrame:
    if not _TEM_PDFPLUMBER:
        raise RuntimeError("pdfplumber não instalado. Execute: pip install pdfplumber")

    # 1. Bordas visíveis
    dados = _tentar_extrair(entrada, "lines")
    if dados:
        df = _normalizar(dados)
        if len(df.columns) > 1: return df

    # 2. Texto agrupado
    dados = _tentar_extrair(entrada, "text")
    if dados:
        df = _normalizar(dados)
        if len(df.columns) > 1: return df

    # 3. Posição X das palavras
    df = _extrair_por_posicao_palavras(entrada)
    if df is not None and len(df.columns) > 1: return df

    # 4. Fallback texto
    return _fallback_texto(entrada)


def _tentar_extrair(entrada: str, estrategia: str) -> list:
    cfg = (
        dict(vertical_strategy="lines", horizontal_strategy="lines",
             snap_tolerance=5, join_tolerance=3)
        if estrategia == "lines" else
        dict(vertical_strategy="text",  horizontal_strategy="text",
             snap_tolerance=5, join_tolerance=5,
             intersection_x_tolerance=15, intersection_y_tolerance=15)
    )
    todas = []
    with pdfplumber.open(entrada) as pdf:
        for page in pdf.pages:
            try:
                for t in page.extract_tables(cfg):
                    if t: todas.extend(t)
            except: pass
    return todas


def _normalizar(linhas: list) -> pd.DataFrame:
    linhas = [r for r in linhas if any(c for c in r if c and str(c).strip())]
    if not linhas: return pd.DataFrame({"Conteúdo": []})
    max_cols = max(len(r) for r in linhas)
    header   = [(str(c).strip() if c else f"Col_{i}") for i, c in enumerate(linhas[0])]
    rows     = [(r + [None]*max_cols)[:max_cols] for r in linhas[1:]]
    df = pd.DataFrame(rows, columns=header)
    df = df.fillna("").map(lambda x: str(x).strip())
    df = df.loc[:, (df != "").any(axis=0)]
    return df


def _extrair_por_posicao_palavras(entrada: str):
    """Detecta colunas pela posição X das palavras — funciona para relatórios sem bordas."""
    todas_linhas = []

    with pdfplumber.open(entrada) as pdf:
        for page in pdf.pages:
            palavras = page.extract_words(x_tolerance=3, y_tolerance=5, keep_blank_chars=False)
            if not palavras: continue
            linha_atual = [palavras[0]]
            for p in palavras[1:]:
                if abs(p['top'] - linha_atual[0]['top']) < 6:
                    linha_atual.append(p)
                else:
                    todas_linhas.append(sorted(linha_atual, key=lambda w: w['x0']))
                    linha_atual = [p]
            if linha_atual:
                todas_linhas.append(sorted(linha_atual, key=lambda w: w['x0']))

    if not todas_linhas: return None

    contagem_gaps = {}
    for linha in todas_linhas:
        for i in range(len(linha) - 1):
            gap = linha[i+1]['x0'] - linha[i]['x1']
            if gap > 10:
                pos = round((linha[i]['x1'] + linha[i+1]['x0']) / 2 / 15) * 15
                contagem_gaps[pos] = contagem_gaps.get(pos, 0) + 1

    if not contagem_gaps: return None

    min_freq     = max(2, len(todas_linhas) * 0.15)
    separadores  = sorted(pos for pos, cnt in contagem_gaps.items() if cnt >= min_freq)
    if not separadores: return None

    seps = [separadores[0]]
    for s in separadores[1:]:
        if s - seps[-1] > 35: seps.append(s)

    limites = [0] + seps + [float('inf')]
    n_cols  = len(limites) - 1
    if n_cols < 2: return None

    matrix = []
    for linha in todas_linhas:
        row = [""] * n_cols
        for p in linha:
            x_mid = (p['x0'] + p['x1']) / 2
            for i in range(n_cols):
                if limites[i] <= x_mid < limites[i+1]:
                    row[i] = (row[i] + " " + p['text']).strip()
                    break
        if any(c for c in row): matrix.append(row)

    if len(matrix) < 2: return None

    header = [(matrix[0][i] if matrix[0][i].strip() else f"Col_{i}") for i in range(n_cols)]
    df = pd.DataFrame(matrix[1:], columns=header)
    df = df.fillna("").map(lambda x: str(x).strip())
    df = df.loc[:, (df != "").any(axis=0)]
    return df


def _fallback_texto(entrada: str) -> pd.DataFrame:
    linhas_raw = []
    with pdfplumber.open(entrada) as pdf:
        for page in pdf.pages:
            txt = page.extract_text(x_tolerance=3, y_tolerance=3)
            if txt:
                for l in txt.split("\n"):
                    if l.strip(): linhas_raw.append(l.strip())
    if not linhas_raw:
        raise ValueError("Nenhum dado extraível. O PDF pode ser uma imagem escaneada.")
    max_cols = 1
    split    = []
    for l in linhas_raw:
        parts = [p.strip() for p in _re.split(r" {2,}", l) if p.strip()]
        split.append(parts)
        if len(parts) > max_cols: max_cols = len(parts)
    if max_cols == 1:
        return pd.DataFrame(linhas_raw, columns=["Conteúdo"])
    header = [(split[0][i] if i < len(split[0]) else f"Col_{i}") for i in range(max_cols)]
    rows   = [(r + [""]*max_cols)[:max_cols] for r in split[1:]]
    return pd.DataFrame(rows, columns=header)


def pdf_para_xlsx(entrada: str, saida: str):
    df = pdf_extrair_dataframe(entrada)
    df.to_excel(saida, index=False, engine="openpyxl")

def pdf_para_csv(entrada: str, saida: str):
    df = pdf_extrair_dataframe(entrada)
    df.to_csv(saida, index=False, encoding="utf-8-sig")


# ─── Imagem → PDF ────────────────────────────────────────────

def imagem_para_pdf(entrada, saida):
    img = Image.open(entrada)
    if img.mode in ("RGBA", "LA", "P"): img = img.convert("RGB")
    img.save(saida, "PDF", resolution=150)


# ─── DOCX/XLSX → Imagem (via PDF) ────────────────────────────

def _pdf_para_imagem_completo(pdf_path: str, saida: str, fmt: str = "png"):
    """Renderiza todas as páginas empilhadas verticalmente."""
    import io
    doc = fitz.open(pdf_path)
    try:
        if doc.page_count == 1:
            pix = doc[0].get_pixmap(matrix=fitz.Matrix(2, 2))
            pix.save(saida)
        else:
            paginas_img = []
            for page in doc:
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")
                paginas_img.append(img)
            larg  = max(i.width  for i in paginas_img)
            alt   = sum(i.height for i in paginas_img) + 10 * (len(paginas_img) - 1)
            tela  = Image.new("RGB", (larg, alt), (240, 240, 240))
            y = 0
            for img in paginas_img:
                tela.paste(img, (0, y)); y += img.height + 10
            if fmt == "jpg": tela.save(saida, "JPEG", quality=92)
            else:            tela.save(saida, "PNG")
    finally:
        doc.close()

def _via_pdf_para_img(entrada: str, saida: str, origem_ext: str, fmt: str):
    saida_dir = os.path.dirname(saida) or "."
    temp_pdf  = os.path.join(saida_dir, f"_tmp_{uuid.uuid4().hex}.pdf")
    try:
        if origem_ext in ("docx", "doc"):   docx_para_pdf(entrada, temp_pdf)
        elif origem_ext in ("xlsx", "xls"): xlsx_para_pdf(entrada, temp_pdf)
        else: raise ValueError(f"Origem '{origem_ext}' não suportada para imagem.")
        _pdf_para_imagem_completo(temp_pdf, saida, fmt=fmt)
    finally:
        if os.path.exists(temp_pdf):
            try: os.remove(temp_pdf)
            except: pass

def docx_para_png(entrada, saida): _via_pdf_para_img(entrada, saida, "docx", "png")
def docx_para_jpg(entrada, saida): _via_pdf_para_img(entrada, saida, "docx", "jpg")
def xlsx_para_png(entrada, saida): _via_pdf_para_img(entrada, saida, "xlsx", "png")
def xlsx_para_jpg(entrada, saida): _via_pdf_para_img(entrada, saida, "xlsx", "jpg")


# ─── Despacho central ─────────────────────────────────────────

_MAPA = {
    ("csv",  "xlsx"): csv_para_xlsx,
    ("xlsx", "csv"):  xlsx_para_csv,
    ("csv",  "pdf"):  csv_para_pdf,
    ("xlsx", "pdf"):  xlsx_para_pdf,
    ("xlsx", "png"):  xlsx_para_png,
    ("xlsx", "jpg"):  xlsx_para_jpg,
    ("pdf",  "docx"): pdf_para_docx,
    ("pdf",  "pptx"): pdf_para_pptx,
    ("pdf",  "ppt"):  pdf_para_ppt,
    ("pdf",  "png"):  pdf_para_png,
    ("pdf",  "xlsx"): pdf_para_xlsx,
    ("pdf",  "csv"):  pdf_para_csv,
    ("docx", "pdf"):  docx_para_pdf,
    ("docx", "png"):  docx_para_png,
    ("docx", "jpg"):  docx_para_jpg,
    ("ppt",  "pdf"):  ppt_para_pdf,
    ("pptx", "pdf"):  ppt_para_pdf,
    ("png",  "pdf"):  imagem_para_pdf,
    ("jpg",  "pdf"):  imagem_para_pdf,
    ("jpeg", "pdf"):  imagem_para_pdf,
}

_ACEITA_ORIENTACAO = {("xlsx","pdf"), ("xls","pdf"), ("csv","pdf")}

def converter_arquivo(entrada, saida, origem, destino, orientacao="retrato"):
    origem  = origem.lower()
    destino = destino.lower()
    if not os.path.exists(entrada):
        raise FileNotFoundError(f"Arquivo não encontrado: {entrada}")
    funcao = _MAPA.get((origem, destino))
    if not funcao:
        raise ValueError(f"Conversão '{origem.upper()}→{destino.upper()}' não suportada.")
    if (origem, destino) in _ACEITA_ORIENTACAO:
        funcao(entrada, saida, orientacao=orientacao)
    else:
        funcao(entrada, saida)